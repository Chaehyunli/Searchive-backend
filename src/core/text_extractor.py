# -*- coding: utf-8 -*-
"""문서 파일에서 텍스트 추출 유틸리티"""
from typing import Optional
from io import BytesIO
import logging

logger = logging.getLogger(__name__)


class TextExtractor:
    """다양한 파일 형식에서 텍스트를 추출하는 유틸리티"""

    @staticmethod
    def extract_text_from_bytes(
        file_data: bytes,
        file_type: str,
        filename: str
    ) -> Optional[str]:
        """
        파일 바이트 데이터에서 텍스트 추출

        Args:
            file_data: 파일 바이트 데이터
            file_type: 파일 MIME 타입
            filename: 파일명

        Returns:
            추출된 텍스트 또는 None
        """
        try:
            if file_type == "application/pdf":
                return TextExtractor._extract_from_pdf(file_data)
            elif file_type == "text/plain":
                return TextExtractor._extract_from_txt(file_data)
            elif file_type in [
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword"
            ]:
                return TextExtractor._extract_from_docx(file_data)
            elif file_type in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel"
            ]:
                return TextExtractor._extract_from_excel(file_data)
            elif file_type in [
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint"
            ]:
                return TextExtractor._extract_from_pptx(file_data)
            elif file_type in [
                "application/x-hwp",
                "application/haansofthwp",
                "application/vnd.hancom.hwp"
            ]:
                return TextExtractor._extract_from_hwp(file_data)
            else:
                logger.warning(f"지원하지 않는 파일 형식: {file_type}")
                return None

        except Exception as e:
            logger.error(f"텍스트 추출 실패: {e}", exc_info=True)
            return None

    @staticmethod # staticmethod 함수는 self의 영향을 받지 않음
    def _extract_from_pdf(file_data: bytes) -> Optional[str]:
        """PDF 파일에서 텍스트 추출"""
        try:
            import pypdf
            pdf_reader = pypdf.PdfReader(BytesIO(file_data))
            text_parts = []

            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

            full_text = "\n".join(text_parts)
            logger.info(f"PDF 텍스트 추출 완료: {len(full_text)} 문자")
            return full_text

        except ImportError:
            logger.error("pypdf 라이브러리가 설치되지 않았습니다. 'pip install pypdf' 실행 필요")
            return None
        except Exception as e:
            logger.error(f"PDF 텍스트 추출 실패: {e}", exc_info=True)
            return None

    @staticmethod
    def _extract_from_txt(file_data: bytes) -> Optional[str]:
        """TXT 파일에서 텍스트 추출"""
        try:
            # UTF-8로 디코딩 시도, 실패 시 CP949(한글) 시도
            try:
                text = file_data.decode('utf-8')
            except UnicodeDecodeError:
                text = file_data.decode('cp949', errors='ignore')

            logger.info(f"TXT 텍스트 추출 완료: {len(text)} 문자")
            return text

        except Exception as e:
            logger.error(f"TXT 텍스트 추출 실패: {e}", exc_info=True)
            return None

    @staticmethod
    def _extract_from_docx(file_data: bytes) -> Optional[str]:
        """DOCX/DOC 파일에서 텍스트 추출"""
        try:
            from docx import Document
            doc = Document(BytesIO(file_data))
            text_parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
            full_text = "\n".join(text_parts)

            logger.info(f"DOCX 텍스트 추출 완료: {len(full_text)} 문자")
            return full_text

        except ImportError:
            logger.error("python-docx 라이브러리가 설치되지 않았습니다. 'pip install python-docx' 실행 필요")
            return None
        except Exception as e:
            logger.error(f"DOCX 텍스트 추출 실패: {e}", exc_info=True)
            return None

    @staticmethod
    def _extract_from_excel(file_data: bytes) -> Optional[str]:
        """Excel 파일에서 텍스트 추출"""
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(BytesIO(file_data), data_only=True)
            text_parts = []

            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text_parts.append(row_text)

            full_text = "\n".join(text_parts)
            logger.info(f"Excel 텍스트 추출 완료: {len(full_text)} 문자")
            return full_text

        except ImportError:
            logger.error("openpyxl 라이브러리가 설치되지 않았습니다. 'pip install openpyxl' 실행 필요")
            return None
        except Exception as e:
            logger.error(f"Excel 텍스트 추출 실패: {e}", exc_info=True)
            return None

    @staticmethod
    def _extract_from_pptx(file_data: bytes) -> Optional[str]:
        """PowerPoint 파일에서 텍스트 추출"""
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(file_data))
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            full_text = "\n".join(text_parts)
            logger.info(f"PPTX 텍스트 추출 완료: {len(full_text)} 문자")
            return full_text

        except ImportError:
            logger.error("python-pptx 라이브러리가 설치되지 않았습니다. 'pip install python-pptx' 실행 필요")
            return None
        except Exception as e:
            logger.error(f"PPTX 텍스트 추출 실패: {e}", exc_info=True)
            return None

    @staticmethod
    def _extract_from_hwp(file_data: bytes) -> Optional[str]:
        """HWP(한글) 파일에서 텍스트 추출 - 개선된 버전"""
        try:
            import olefile
            import zlib
            import struct

            ole = olefile.OleFileIO(BytesIO(file_data))
            text_parts = []

            # HWP 5.0 이상 버전 처리
            if ole.exists('BodyText'):
                # 모든 Section 스트림 찾기
                sections = []
                for entry in ole.listdir():
                    if len(entry) == 2 and entry[0] == 'BodyText' and entry[1].startswith('Section'):
                        sections.append(entry)

                # Section 순서대로 정렬 (Section0, Section1, ...)
                sections.sort(key=lambda x: int(x[1].replace('Section', '')))

                for section in sections:
                    try:
                        stream = ole.openstream(section)
                        data = stream.read()

                        if not data:
                            continue

                        # 압축 해제 시도
                        try:
                            # HWP는 zlib 압축을 사용
                            unpacked = zlib.decompress(data, -15)
                        except:
                            # 압축되지 않은 경우
                            unpacked = data

                        # HWP 레코드 구조 파싱
                        text = TextExtractor._parse_hwp_text(unpacked)
                        if text and text.strip():
                            text_parts.append(text)

                    except Exception as section_error:
                        logger.debug(f"Section {section} 처리 중 오류: {section_error}")
                        continue

            ole.close()

            if text_parts:
                full_text = "\n".join(text_parts)
                # 추가 정리: 연속된 공백 제거
                import re
                full_text = re.sub(r'\s+', ' ', full_text)
                full_text = full_text.strip()

                logger.info(f"HWP 텍스트 추출 완료: {len(full_text)} 문자")
                return full_text
            else:
                logger.warning("HWP 파일에서 텍스트를 추출할 수 없습니다.")
                return None

        except ImportError:
            logger.error("olefile 라이브러리가 설치되지 않았습니다. 'pip install olefile' 실행 필요")
            return None
        except Exception as e:
            logger.error(f"HWP 텍스트 추출 실패: {e}", exc_info=True)
            return None

    @staticmethod
    def _parse_hwp_text(data: bytes) -> str:
        """HWP 바이너리 데이터에서 텍스트 파싱"""
        import struct

        text_parts = []
        pos = 0

        try:
            while pos < len(data):
                # 최소 레코드 헤더 크기 확인
                if pos + 4 > len(data):
                    break

                # 레코드 헤더 읽기 (4바이트)
                try:
                    record_header = struct.unpack('<I', data[pos:pos+4])[0]
                except:
                    pos += 1
                    continue

                record_id = record_header & 0x3FF
                record_level = (record_header >> 10) & 0x3FF
                record_size = (record_header >> 20) & 0xFFF

                pos += 4

                # 레코드 크기가 0xFFF이면 추가로 4바이트 읽기
                if record_size == 0xFFF:
                    if pos + 4 > len(data):
                        break
                    try:
                        record_size = struct.unpack('<I', data[pos:pos+4])[0]
                        pos += 4
                    except:
                        break

                # 레코드 데이터 영역
                if pos + record_size > len(data):
                    break

                record_data = data[pos:pos+record_size]
                pos += record_size

                # 텍스트 레코드 (HWPTAG_PARA_TEXT = 67)
                if record_id == 67:
                    try:
                        # UTF-16LE로 디코딩
                        text = record_data.decode('utf-16le', errors='ignore')
                        # NULL 문자 제거
                        text = text.replace('\x00', '')
                        # 제어 문자 제거 (줄바꿈과 탭은 유지)
                        cleaned_text = ''.join(
                            char for char in text
                            if char.isprintable() or char in ['\n', '\r', '\t', ' ']
                        )
                        if cleaned_text.strip():
                            text_parts.append(cleaned_text)
                    except Exception as e:
                        logger.debug(f"텍스트 디코딩 오류: {e}")
                        continue

        except Exception as e:
            logger.debug(f"HWP 파싱 오류: {e}")

        return ' '.join(text_parts)


# 전역 텍스트 추출기 인스턴스
text_extractor = TextExtractor()
